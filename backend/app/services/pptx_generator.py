"""
PPTX演示文稿生成服务

基于python-pptx的PowerPoint生成
"""
from typing import Optional, List, Dict, Any
from datetime import datetime
from io import BytesIO
import logging
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

from app.config import settings

logger = logging.getLogger("app")


# 颜色主题 - 青色主题
class ThemeColors:
    """主题色彩"""
    PRIMARY = RGBColor(0x08, 0x91, 0xB2)      # #0891B2
    PRIMARY_LIGHT = RGBColor(0x22, 0xD3, 0xEE)  # #22D3EE
    PRIMARY_DARK = RGBColor(0x0E, 0x74, 0x90)   # #0E7490
    SUCCESS = RGBColor(0x05, 0x96, 0x69)       # #059669
    TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)     # #1E293B
    TEXT_BODY = RGBColor(0x47, 0x55, 0x69)      # #475569
    TEXT_SECONDARY = RGBColor(0x94, 0xA3, 0xB8) # #94A3B8
    BG_LIGHT = RGBColor(0xF8, 0xFA, 0xFC)        # #F8FAFC
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)


class PptxGeneratorService:
    """PPTX演示文稿生成服务"""

    # 幻灯片尺寸 (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # 字体设置
    FONT_NAME_CN = "微软雅黑"
    FONT_NAME_EN = "Arial"

    def __init__(self):
        self.logger = logging.getLogger("app")

    def _is_template_content(self, text: str) -> bool:
        """检测是否为模板占位符内容"""
        import re
        if not text:
            return True
        template_phrases = [
            "详细描述讨论了什么问题",
            "各方观点",
            "最终结论是什么",
            "决定的具体内容",
            "任务内容说明",
            "负责人：XXX",
            "截止时间：XXX",
            "根据转写内容分析",
            "根据转写内容总结",
            "根据转写内容",
        ]
        text_lower = text.lower()
        match_count = sum(1 for phrase in template_phrases if phrase.lower() in text_lower)
        return match_count >= 2

    def _filter_list_field(self, items) -> list:
        """过滤列表字段，移除模板内容"""
        import re
        if not items:
            return []
        if isinstance(items, str):
            items = [items]
        if not isinstance(items, list):
            return [str(items)]

        result = []
        for item in items:
            item_str = str(item).strip()
            if not item_str or len(item_str) < 10:
                continue
            # 清理代码格式
            item_str = re.sub(r"\{[^}]*\}", lambda m: m.group(0).replace("{", "").replace("}", "").replace("'", ""), item_str)
            item_str = item_str.replace("{", "").replace("}", "").replace("'", "").strip()
            # 过滤模板内容
            if item_str and not self._is_template_content(item_str) and len(item_str) > 10:
                result.append(item_str)
        return result

    async def generate_course_pptx(self, summary: dict) -> bytes:
        """
        生成课程总结PPT

        Args:
            summary: 课程总结数据字典，包含以下字段:
                - title: 标题
                - topic: 主题
                - key_points: 核心知识点（字符串，多行或用|分隔）
                - difficulties: 重点难点
                - memorable_quote: 金句
                - next_suggestion: 预习建议
                - date: 可选日期

        Returns:
            PPTX文件字节数据
        """
        import re

        self.logger.info(f"[PPTX] generate_course_pptx started, title: {summary.get('title', 'Untitled')}")

        # 预处理：清理代码格式
        def clean_text(text):
            if not text:
                return text
            text = re.sub(r"\{[^}]*\}", lambda m: m.group(0).replace("{", "").replace("}", "").replace("'", ""), str(text))
            text = text.replace("{", "").replace("}", "").replace("'", "").strip()
            return text

        for key in ['key_points', 'difficulties', 'memorable_quote', 'next_suggestion', 'full_text']:
            if key in summary and summary[key]:
                summary[key] = clean_text(summary[key])

        try:
            # 创建演示文稿
            prs = Presentation()
            prs.slide_width = self.SLIDE_WIDTH
            prs.slide_height = self.SLIDE_HEIGHT

            # 生成各页幻灯片
            self._create_cover_slide(prs, summary)
            self._create_topic_slide(prs, summary)
            self._create_key_points_slide(prs, summary)
            self._create_difficulties_slide(prs, summary)
            self._create_quotes_slide(prs, summary)
            self._create_suggestion_slide(prs, summary)
            # 添加完整总结页
            self._create_full_summary_slide(prs, summary)
            self._create_ending_slide(prs, summary)

            slide_count = len(prs.slides)
            self.logger.info(f"[PPTX] Generated {slide_count} slides successfully")

            # 保存到字节流
            output = BytesIO()
            prs.save(output)
            output.seek(0)

            self.logger.info(f"[PPTX] PPTX generation completed, size: {output.getbuffer().nbytes} bytes")
            return output.getvalue()

        except Exception as e:
            self.logger.error(f"[PPTX] Generation failed: {str(e)}", exc_info=True)
            raise

    def _create_cover_slide(self, prs: Presentation, summary: dict):
        """创建封面页"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加背景装饰 - 顶部渐变条
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY
        shape.line.fill.background()

        # 添加装饰圆形
        circle1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(10), Inches(1.5),
            Inches(2), Inches(2)
        )
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = ThemeColors.PRIMARY_LIGHT
        circle1.line.fill.background()

        circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(11.5), Inches(0.5),
            Inches(1.2), Inches(1.2)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = ThemeColors.WHITE
        circle2.line.fill.background()

        # 标题
        title = summary.get("title", "课程总结")
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.TEXT_DARK
        p.font.name = self.FONT_NAME_CN
        p.alignment = PP_ALIGN.LEFT

        # 副标题
        topic = summary.get("topic", "")
        if topic:
            topic_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11), Inches(0.8))
            tf = topic_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"主题：{topic}"
            p.font.size = Pt(24)
            p.font.color.rgb = ThemeColors.PRIMARY
            p.font.name = self.FONT_NAME_CN
            p.alignment = PP_ALIGN.LEFT

        # 日期
        date_str = summary.get("date", datetime.now().strftime("%Y年%m月%d日"))
        date_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = date_str
        p.font.size = Pt(18)
        p.font.color.rgb = ThemeColors.TEXT_SECONDARY
        p.font.name = self.FONT_NAME_CN
        p.alignment = PP_ALIGN.LEFT

        self.logger.debug(f"[PPTX] Cover slide created: {title}")

    def _create_topic_slide(self, prs: Presentation, summary: dict):
        """创建课程主题页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 顶部装饰条
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "课程主题"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.PRIMARY
        p.font.name = self.FONT_NAME_CN

        # 内容区域 - 带背景的卡片
        content_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(2),
            Inches(11.5), Inches(4.5)
        )
        content_shape.fill.solid()
        content_shape.fill.fore_color.rgb = ThemeColors.BG_LIGHT
        content_shape.line.color.rgb = ThemeColors.PRIMARY
        content_shape.line.width = Pt(2)

        # 主题文字
        topic = summary.get("topic", "本节课主题")
        topic_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10), Inches(2.5))
        tf = topic_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = topic
        p.font.size = Pt(28)
        p.font.color.rgb = ThemeColors.TEXT_DARK
        p.font.name = self.FONT_NAME_CN
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(20)

        self.logger.debug(f"[PPTX] Topic slide created")

    def _create_key_points_slide(self, prs: Presentation, summary: dict):
        """创建核心知识点页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 顶部装饰条
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "核心知识点"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.PRIMARY
        p.font.name = self.FONT_NAME_CN

        # 解析知识点 - 支持多种格式
        key_points_str = summary.get("key_points", "")
        if isinstance(key_points_str, list):
            key_points = [str(kp) for kp in key_points_str]
        elif "|" in key_points_str:
            key_points = [kp.strip() for kp in key_points_str.split("|") if kp.strip()]
        else:
            key_points = [kp.strip() for kp in key_points_str.split("\n") if kp.strip()]

        if not key_points:
            key_points = ["暂无核心知识点"]

        # 根据知识点数量决定布局
        if len(key_points) <= 3:
            # 3个以内，每个占较大空间
            card_height = Inches(1.4)
            start_y = 2.0
            for i, point in enumerate(key_points):
                y_pos = start_y + i * (card_height + Inches(0.2))

                # 序号圆圈
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(0.8), y_pos + Inches(0.35),
                    Inches(0.7), Inches(0.7)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = ThemeColors.PRIMARY
                circle.line.fill.background()

                num_frame = circle.text_frame
                num_frame.paragraphs[0].text = str(i + 1)
                num_frame.paragraphs[0].font.size = Pt(22)
                num_frame.paragraphs[0].font.bold = True
                num_frame.paragraphs[0].font.color.rgb = ThemeColors.WHITE
                num_frame.paragraphs[0].font.name = self.FONT_NAME_EN
                num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                # 内容卡片
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(1.8), y_pos,
                    Inches(10.5), card_height
                )
                card.fill.solid()
                card.fill.fore_color.rgb = ThemeColors.BG_LIGHT
                card.line.color.rgb = ThemeColors.PRIMARY_LIGHT
                card.line.width = Pt(1.5)

                text_box = slide.shapes.add_textbox(Inches(2.1), y_pos + Inches(0.2), Inches(10), card_height - Inches(0.3))
                tf = text_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = point
                p.font.size = Pt(22)
                p.font.color.rgb = ThemeColors.TEXT_DARK
                p.font.name = self.FONT_NAME_CN
                p.line_spacing = 1.3
        else:
            # 4个以上，压缩布局
            card_height = Inches(0.75)
            start_y = 1.8
            card_spacing = Inches(0.1)

            for i, point in enumerate(key_points[:6]):
                y_pos = start_y + i * (card_height + card_spacing)

                # 序号圆圈
                circle = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    Inches(0.8), y_pos + Inches(0.08),
                    Inches(0.55), Inches(0.55)
                )
                circle.fill.solid()
                circle.fill.fore_color.rgb = ThemeColors.PRIMARY
                circle.line.fill.background()

                num_frame = circle.text_frame
                num_frame.paragraphs[0].text = str(i + 1)
                num_frame.paragraphs[0].font.size = Pt(16)
                num_frame.paragraphs[0].font.bold = True
                num_frame.paragraphs[0].font.color.rgb = ThemeColors.WHITE
                num_frame.paragraphs[0].font.name = self.FONT_NAME_EN
                num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                # 知识点内容卡片
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(1.55), y_pos,
                    Inches(10.75), card_height
                )
                card.fill.solid()
                card.fill.fore_color.rgb = ThemeColors.BG_LIGHT
                card.line.color.rgb = ThemeColors.PRIMARY_LIGHT
                card.line.width = Pt(1)

                # 知识点文字
                text_box = slide.shapes.add_textbox(Inches(1.8), y_pos + Inches(0.1), Inches(10.3), card_height)
                tf = text_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = point
                p.font.size = Pt(18)
                p.font.color.rgb = ThemeColors.TEXT_DARK
                p.font.name = self.FONT_NAME_CN

        self.logger.debug(f"[PPTX] Key points slide created with {len(key_points)} points")

    def _create_difficulties_slide(self, prs: Presentation, summary: dict):
        """创建重点难点页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 顶部装饰条
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.SUCCESS
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "重点难点"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.SUCCESS
        p.font.name = self.FONT_NAME_CN

        # 解析难点
        difficulties_str = summary.get("difficulties", "")
        if isinstance(difficulties_str, list):
            difficulties = [str(d) for d in difficulties_str]
        elif "|" in difficulties_str:
            difficulties = [d.strip() for d in difficulties_str.split("|") if d.strip()]
        else:
            difficulties = [d.strip() for d in difficulties_str.split("\n") if d.strip()]

        if not difficulties:
            difficulties = ["暂无重点难点"]

        # 根据数量调整布局
        if len(difficulties) <= 2:
            col_width = Inches(11)
            start_y = Inches(2)
            card_height = Inches(1.8)
            for i, difficulty in enumerate(difficulties):
                y_pos = start_y + i * (card_height + Inches(0.3))

                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(0.8), y_pos,
                    col_width, card_height
                )
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xF0, 0xF9, 0xF4)
                card.line.color.rgb = ThemeColors.SUCCESS
                card.line.width = Pt(2)

                # 内容
                content_box = slide.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.3), col_width - Inches(0.8), card_height - Inches(0.5))
                tf = content_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = difficulty
                p.font.size = Pt(22)
                p.font.color.rgb = ThemeColors.TEXT_DARK
                p.font.name = self.FONT_NAME_CN
                p.line_spacing = 1.4
        else:
            # 左右两栏布局
            col_width = Inches(5.5)
            start_y = Inches(2)
            card_height = Inches(1.2)

            for i, difficulty in enumerate(difficulties[:4]):
                col = i % 2
                row = i // 2

                x_pos = Inches(0.8) + col * (col_width + Inches(0.5))
                y_pos = start_y + row * (card_height + Inches(0.3))

                # 难点卡片
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    x_pos, y_pos,
                    col_width, card_height
                )
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xF0, 0xF9, 0xF4)
                card.line.color.rgb = ThemeColors.SUCCESS
                card.line.width = Pt(1.5)

                # 内容
                content_box = slide.shapes.add_textbox(x_pos + Inches(0.3), y_pos + Inches(0.2), col_width - Inches(0.5), card_height - Inches(0.3))
                tf = content_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = difficulty
                p.font.size = Pt(18)
                p.font.color.rgb = ThemeColors.TEXT_DARK
                p.font.name = self.FONT_NAME_CN
                p.line_spacing = 1.3

        self.logger.debug(f"[PPTX] Difficulties slide created with {len(difficulties)} items")

    def _create_quotes_slide(self, prs: Presentation, summary: dict):
        """创建金句页"""
        import re
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 顶部装饰条
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY_DARK
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "金句摘录"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.PRIMARY_DARK
        p.font.name = self.FONT_NAME_CN

        # 金句内容 - 清理代码格式
        quote = summary.get("memorable_quote", "暂无金句")
        # 清理代码格式
        quote = re.sub(r"\{[^}]*\}", lambda m: m.group(0).replace("{", "").replace("}", "").replace("'", ""), str(quote))
        quote = quote.replace("{", "").replace("}", "").replace("'", "").strip()

        # 大引号装饰
        left_quote = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(1), Inches(1))
        tf = left_quote.text_frame
        p = tf.paragraphs[0]
        p.text = """
        p.font.size = Pt(72)
        p.font.color.rgb = ThemeColors.PRIMARY_LIGHT
        p.font.name = "Georgia"

        # 引用内容卡片
        content_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), Inches(2.5),
            Inches(11), Inches(3.5)
        )
        content_card.fill.solid()
        content_card.fill.fore_color.rgb = ThemeColors.BG_LIGHT
        content_card.line.color.rgb = ThemeColors.PRIMARY
        content_card.line.width = Pt(2)

        # 引用文字
        quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10), Inches(2.5))
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = quote if quote else "暂无金句"
        p.font.size = Pt(26)
        p.font.color.rgb = ThemeColors.TEXT_DARK
        p.font.name = self.FONT_NAME_CN
        p.alignment = PP_ALIGN.CENTER

        # 右引号
        right_quote = slide.shapes.add_textbox(Inches(11.5), Inches(5), Inches(1), Inches(1))
        tf = right_quote.text_frame
        p = tf.paragraphs[0]
        p.text = """
        p.font.size = Pt(72)
        p.font.color.rgb = ThemeColors.PRIMARY_LIGHT
        p.font.name = "Georgia"

        self.logger.debug(f"[PPTX] Quotes slide created")

    def _create_suggestion_slide(self, prs: Presentation, summary: dict):
        """创建预习建议页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 顶部装饰条
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "预习建议"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.PRIMARY
        p.font.name = self.FONT_NAME_CN

        # 建议内容
        suggestion = summary.get("next_suggestion", "暂无预习建议")

        # 建议卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8), Inches(2),
            Inches(11.5), Inches(4.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = ThemeColors.BG_LIGHT
        card.line.color.rgb = ThemeColors.PRIMARY
        card.line.width = Pt(2)

        # 图标区域 - 书本图标用形状模拟
        icon_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.3), Inches(2.8),
            Inches(1), Inches(1)
        )
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = ThemeColors.PRIMARY
        icon_shape.line.fill.background()

        # 建议文字
        content_box = slide.shapes.add_textbox(Inches(2.8), Inches(2.8), Inches(9), Inches(3.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = suggestion
        p.font.size = Pt(24)
        p.font.color.rgb = ThemeColors.TEXT_DARK
        p.font.name = self.FONT_NAME_CN
        p.alignment = PP_ALIGN.LEFT

        self.logger.debug(f"[PPTX] Suggestion slide created")

    def _create_full_summary_slide(self, prs: Presentation, summary: dict):
        """创建完整总结页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 顶部装饰条
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY_DARK
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "课程总结"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.PRIMARY_DARK
        p.font.name = self.FONT_NAME_CN

        # 完整总结内容
        full_text = summary.get("full_text", "")

        if full_text:
            # 创建内容区域
            content_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(2),
                Inches(11.5), Inches(4.8)
            )
            content_card.fill.solid()
            content_card.fill.fore_color.rgb = ThemeColors.BG_LIGHT
            content_card.line.color.rgb = ThemeColors.PRIMARY
            content_card.line.width = Pt(2)

            # 内容文字
            content_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(10.7), Inches(4.3))
            tf = content_box.text_frame
            tf.word_wrap = True

            # 分段落显示
            paragraphs = str(full_text).split('\n')
            for i, para_text in enumerate(paragraphs):
                if para_text.strip():
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = para_text.strip()
                    p.font.size = Pt(20)
                    p.font.color.rgb = ThemeColors.TEXT_DARK
                    p.font.name = self.FONT_NAME_CN
                    p.line_spacing = 1.5
                    p.space_after = Pt(12)
        else:
            # 无内容提示
            content_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
            tf = content_box.text_frame
            p = tf.paragraphs[0]
            p.text = "完整总结内容整理中..."
            p.font.size = Pt(24)
            p.font.color.rgb = ThemeColors.TEXT_SECONDARY
            p.font.name = self.FONT_NAME_CN
            p.alignment = PP_ALIGN.CENTER

        self.logger.debug(f"[PPTX] Full summary slide created")

    def _create_ending_slide(self, prs: Presentation, summary: dict):
        """创建结束页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 全屏背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = ThemeColors.PRIMARY
        bg.line.fill.background()

        # 装饰圆形
        circle1 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-1), Inches(-1),
            Inches(4), Inches(4)
        )
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = ThemeColors.PRIMARY_LIGHT
        circle1.line.fill.background()

        circle2 = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(10), Inches(5),
            Inches(4), Inches(4)
        )
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = ThemeColors.PRIMARY_DARK
        circle2.line.fill.background()

        # 结束语
        ending_box = slide.shapes.add_textbox(Inches(0), Inches(2.8), self.SLIDE_WIDTH, Inches(1.5))
        tf = ending_box.text_frame
        p = tf.paragraphs[0]
        p.text = "感谢观看"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.WHITE
        p.font.name = self.FONT_NAME_CN
        p.alignment = PP_ALIGN.CENTER

        # 副标题
        sub_box = slide.shapes.add_textbox(Inches(0), Inches(4.5), self.SLIDE_WIDTH, Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = "AI小商 · 智慧校园助手"
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0xCC, 0xEE, 0xFF)
        p.font.name = self.FONT_NAME_CN
        p.alignment = PP_ALIGN.CENTER

        self.logger.debug(f"[PPTX] Ending slide created")

    async def generate_meeting_pptx(self, summary: dict) -> bytes:
        """
        生成会议总结PPT

        Args:
            summary: 会议总结数据字典，包含以下字段:
                - title: 标题
                - topic: 主题
                - discussion_points: 讨论要点
                - resolutions: 决议事项
                - action_items: 行动项
                - memorable_quote: 重要发言
                - full_text: 完整总结
                - date: 可选日期

        Returns:
            PPTX文件字节数据
        """
        import re

        self.logger.info(f"[PPTX] generate_meeting_pptx started, title: {summary.get('title', 'Untitled')}")

        # 预处理：清理代码格式
        def clean_text(text):
            if not text:
                return text
            text = re.sub(r"\{[^}]*\}", lambda m: m.group(0).replace("{", "").replace("}", "").replace("'", ""), str(text))
            text = text.replace("{", "").replace("}", "").replace("'", "").strip()
            return text

        for key in ['discussion_points', 'resolutions', 'action_items', 'memorable_quote', 'full_text']:
            if key in summary and summary[key]:
                summary[key] = clean_text(summary[key])

        # 过滤模板内容
        summary['discussion_points'] = self._filter_list_field(summary.get('discussion_points', ''))
        summary['resolutions'] = self._filter_list_field(summary.get('resolutions', ''))
        summary['action_items'] = self._filter_list_field(summary.get('action_items', ''))

        try:
            # 创建演示文稿
            prs = Presentation()
            prs.slide_width = self.SLIDE_WIDTH
            prs.slide_height = self.SLIDE_HEIGHT

            # 生成各页幻灯片
            self._create_meeting_cover_slide(prs, summary)
            self._create_meeting_topic_slide(prs, summary)
            self._create_discussion_points_slide(prs, summary)
            self._create_resolutions_slide(prs, summary)
            self._create_action_items_slide(prs, summary)
            self._create_meeting_quote_slide(prs, summary)
            self._create_meeting_full_text_slide(prs, summary)
            self._create_ending_slide(prs, summary)

            slide_count = len(prs.slides)
            self.logger.info(f"[PPTX] Generated meeting PPT {slide_count} slides successfully")

            # 保存到字节流
            output = BytesIO()
            prs.save(output)
            output.seek(0)

            self.logger.info(f"[PPTX] Meeting PPTX generation completed, size: {output.getbuffer().nbytes} bytes")
            return output.getvalue()

        except Exception as e:
            self.logger.error(f"[PPTX] Meeting PPTX generation failed: {str(e)}", exc_info=True)
            raise

    def _create_meeting_cover_slide(self, prs: Presentation, summary: dict):
        """创建会议封面页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 顶部装饰条
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY
        shape.line.fill.background()

        # 装饰圆形
        circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(1.5), Inches(2), Inches(2))
        circle1.fill.solid()
        circle1.fill.fore_color.rgb = ThemeColors.PRIMARY_LIGHT
        circle1.line.fill.background()

        circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.5), Inches(0.5), Inches(1.2), Inches(1.2))
        circle2.fill.solid()
        circle2.fill.fore_color.rgb = ThemeColors.WHITE
        circle2.line.fill.background()

        # 标题
        title = summary.get("title", "会议总结")
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.TEXT_DARK
        p.font.name = self.FONT_NAME_CN
        p.alignment = PP_ALIGN.LEFT

        # 副标题
        topic = summary.get("topic", "")
        if topic:
            topic_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(11), Inches(0.8))
            tf = topic_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"主题：{topic}"
            p.font.size = Pt(28)
            p.font.color.rgb = ThemeColors.TEXT_BODY
            p.font.name = self.FONT_NAME_CN

        # 日期
        date = summary.get("date", "")
        if date:
            date_box = slide.shapes.add_textbox(Inches(1), Inches(5.6), Inches(11), Inches(0.5))
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = date
            p.font.size = Pt(16)
            p.font.color.rgb = ThemeColors.TEXT_SECONDARY
            p.font.name = self.FONT_NAME_CN

    def _create_meeting_topic_slide(self, prs: Presentation, summary: dict):
        """创建会议主题页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 左侧色块
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.3), self.SLIDE_HEIGHT
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "会议主题"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.PRIMARY
        p.font.name = self.FONT_NAME_CN

        # 主题内容
        topic = summary.get("topic", "")
        if topic:
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(3))
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = topic
            p.font.size = Pt(28)
            p.font.color.rgb = ThemeColors.TEXT_DARK
            p.font.name = self.FONT_NAME_CN
            p.line_spacing = 1.5

    def _create_discussion_points_slide(self, prs: Presentation, summary: dict):
        """创建讨论要点页"""
        import re
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 左侧色块
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), self.SLIDE_HEIGHT)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "一、讨论要点"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.PRIMARY
        p.font.name = self.FONT_NAME_CN

        # 解析内容
        discussion_points = summary.get("discussion_points", "")
        if not discussion_points:
            discussion_points = summary.get("key_points", "")

        # 处理多种格式
        if isinstance(discussion_points, list):
            points = [str(p) for p in discussion_points]
        elif "|" in str(discussion_points):
            points = [p.strip() for p in str(discussion_points).split("|") if p.strip()]
        else:
            points = [p.strip() for p in str(discussion_points).split("\n") if p.strip()]

        # 清理代码格式
        cleaned_points = []
        for point in points:
            # 移除代码格式
            point = re.sub(r"\{[^}]+\}", lambda m: m.group().replace("{", "").replace("}", "").replace("'", ""), point)
            point = point.strip()
            if point and len(point) > 3:
                cleaned_points.append(point)

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        if cleaned_points:
            for i, point in enumerate(cleaned_points):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {point}"
                p.font.size = Pt(20)
                p.font.color.rgb = ThemeColors.TEXT_DARK
                p.font.name = self.FONT_NAME_CN
                p.space_after = Pt(14)
                p.line_spacing = 1.3
        else:
            p = tf.paragraphs[0]
            p.text = "暂无讨论要点"
            p.font.size = Pt(20)
            p.font.color.rgb = ThemeColors.TEXT_SECONDARY

    def _create_resolutions_slide(self, prs: Presentation, summary: dict):
        """创建决议事项页"""
        import re
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 左侧色块
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), self.SLIDE_HEIGHT)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.SUCCESS
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "二、决议事项"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.SUCCESS
        p.font.name = self.FONT_NAME_CN

        # 解析内容
        resolutions = summary.get("resolutions", "")
        if not resolutions:
            resolutions = summary.get("difficulties", "")

        if isinstance(resolutions, list):
            items = [str(r) for r in resolutions]
        elif "|" in str(resolutions):
            items = [r.strip() for r in str(resolutions).split("|") if r.strip()]
        else:
            items = [r.strip() for r in str(resolutions).split("\n") if r.strip()]

        # 清理代码格式
        cleaned_items = []
        for item in items:
            item = re.sub(r"\{[^}]+\}", lambda m: m.group().replace("{", "").replace("}", "").replace("'", ""), item)
            item = item.strip()
            if item and len(item) > 3:
                cleaned_items.append(item)

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        if cleaned_items:
            for i, item in enumerate(cleaned_items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"✓ {item}"
                p.font.size = Pt(22)
                p.font.color.rgb = ThemeColors.TEXT_DARK
                p.font.name = self.FONT_NAME_CN
                p.space_after = Pt(18)
                p.line_spacing = 1.3
        else:
            p = tf.paragraphs[0]
            p.text = "暂无决议事项"
            p.font.size = Pt(20)
            p.font.color.rgb = ThemeColors.TEXT_SECONDARY

    def _create_action_items_slide(self, prs: Presentation, summary: dict):
        """创建待办事项页"""
        import re
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 左侧色块
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), self.SLIDE_HEIGHT)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xF5, 0x9E, 0x0B)  # 橙色
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "三、待办事项"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xF5, 0x9E, 0x0B)
        p.font.name = self.FONT_NAME_CN

        # 解析内容
        action_items = summary.get("action_items", "")
        if not action_items:
            action_items = summary.get("next_suggestion", "")

        if isinstance(action_items, list):
            items = [str(a) for a in action_items]
        elif "|" in str(action_items):
            items = [a.strip() for a in str(action_items).split("|") if a.strip()]
        else:
            items = [a.strip() for a in str(action_items).split("\n") if a.strip()]

        # 清理代码格式
        cleaned_items = []
        for item in items:
            item = re.sub(r"\{[^}]+\}", lambda m: m.group().replace("{", "").replace("}", "").replace("'", ""), item)
            item = item.strip()
            if item and len(item) > 3:
                cleaned_items.append(item)

        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        if cleaned_items:
            for i, item in enumerate(cleaned_items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"▸ {item}"
                p.font.size = Pt(22)
                p.font.color.rgb = ThemeColors.TEXT_DARK
                p.font.name = self.FONT_NAME_CN
                p.space_after = Pt(18)
                p.line_spacing = 1.3
        else:
            p = tf.paragraphs[0]
            p.text = "暂无待办事项"
            p.font.size = Pt(20)
            p.font.color.rgb = ThemeColors.TEXT_SECONDARY

    def _create_meeting_quote_slide(self, prs: Presentation, summary: dict):
        """创建重要发言页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 左侧色块
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), self.SLIDE_HEIGHT)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "重要发言"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.PRIMARY
        p.font.name = self.FONT_NAME_CN

        # 内容
        quote = summary.get("memorable_quote", "")
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11.5), Inches(4))
        tf = content_box.text_frame
        tf.word_wrap = True

        if quote:
            p = tf.paragraphs[0]
            p.text = f'"{quote}"'
            p.font.size = Pt(26)
            p.font.color.rgb = ThemeColors.TEXT_DARK
            p.font.name = self.FONT_NAME_CN
            p.font.italic = True
        else:
            p = tf.paragraphs[0]
            p.text = "暂无重要发言记录"
            p.font.size = Pt(20)
            p.font.color.rgb = ThemeColors.TEXT_SECONDARY

    def _create_meeting_full_text_slide(self, prs: Presentation, summary: dict):
        """创建会议完整记录页"""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # 左侧色块
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), self.SLIDE_HEIGHT)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ThemeColors.PRIMARY_DARK
        shape.line.fill.background()

        # 标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "会议纪要"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = ThemeColors.PRIMARY_DARK
        p.font.name = self.FONT_NAME_CN

        # 完整记录内容
        full_text = summary.get("full_text", "")

        if full_text:
            # 内容区域
            content_card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(1.6),
                Inches(11.5), Inches(5.2)
            )
            content_card.fill.solid()
            content_card.fill.fore_color.rgb = ThemeColors.BG_LIGHT
            content_card.line.color.rgb = ThemeColors.PRIMARY
            content_card.line.width = Pt(2)

            # 内容文字
            content_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(10.7), Inches(4.6))
            tf = content_box.text_frame
            tf.word_wrap = True

            # 分段落显示
            paragraphs = str(full_text).split('\n')
            for i, para_text in enumerate(paragraphs):
                if para_text.strip():
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = para_text.strip()
                    p.font.size = Pt(20)
                    p.font.color.rgb = ThemeColors.TEXT_DARK
                    p.font.name = self.FONT_NAME_CN
                    p.line_spacing = 1.5
                    p.space_after = Pt(12)
        else:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(1))
            tf = content_box.text_frame
            p = tf.paragraphs[0]
            p.text = "会议完整记录整理中..."
            p.font.size = Pt(24)
            p.font.color.rgb = ThemeColors.TEXT_SECONDARY
            p.font.name = self.FONT_NAME_CN
            p.alignment = PP_ALIGN.CENTER

        self.logger.debug(f"[PPTX] Meeting full text slide created")

    async def save_to_file(self, content: bytes, filename: str) -> str:
        """
        保存演示文稿到文件

        Args:
            content: 文件内容
            filename: 文件名

        Returns:
            保存后的文件路径
        """
        save_dir = Path(settings.UPLOAD_DIR) / "presentations"
        save_dir.mkdir(parents=True, exist_ok=True)

        filepath = save_dir / filename
        with open(filepath, "wb") as f:
            f.write(content)

        self.logger.info(f"[PPTX] File saved to: {filepath}")
        return str(filepath)


    async def generate_review_summary(
        self,
        title: str,
        topic: str,
        key_points: str,
        difficulties: str,
        memorable_quote: str,
        next_suggestion: str,
        full_text: str,
        metadata: dict = None
    ) -> bytes:
        """
        生成录音回顾总结PPT

        Args:
            title: 标题
            topic: 主题
            key_points: 核心知识点（课程用）或讨论要点（会议用）
            difficulties: 重点难点（课程用）或决议事项（会议用）
            memorable_quote: 金句/重要发言
            next_suggestion: 预习建议（课程用）或待办事项（会议用）
            full_text: 完整总结
            metadata: 额外元数据 (record_type: course/meeting)

        Returns:
            PPTX文件字节数据
        """
        record_type = metadata.get("record_type", "course") if metadata else "course"

        summary = {
            "title": title,
            "topic": topic,
            "key_points": key_points,
            "difficulties": difficulties,
            "memorable_quote": memorable_quote,
            "next_suggestion": next_suggestion,
            "full_text": full_text,
            "date": metadata.get("created_at", "") if metadata else ""
        }

        if record_type == "meeting":
            # 会议类型使用会议专用字段
            summary["discussion_points"] = key_points
            summary["resolutions"] = difficulties
            summary["action_items"] = next_suggestion
            return await self.generate_meeting_pptx(summary)
        else:
            return await self.generate_course_pptx(summary)


# 单例实例
pptx_generator = PptxGeneratorService()